"""
create_how_it_works_presentation.py

Generates docs/content-dashboard-how-it-works.pptx — a technical walkthrough
of the Content Performance Dashboard pipeline plus Method A/B/C comparison.

Usage:
    python docs/create_how_it_works_presentation.py
"""

import json
import subprocess
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (copied verbatim from create_presentation.py) ─────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

BLUE_DARK  = rgb(0x00, 0x53, 0xA0)
BLUE_MID   = rgb(0x00, 0x78, 0xD4)
GREEN      = rgb(0x10, 0x79, 0x47)
ORANGE     = rgb(0xD8, 0x34, 0x00)
WHITE      = rgb(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = rgb(0xF3, 0xF2, 0xF1)
MID_GRAY   = rgb(0x60, 0x5E, 0x5C)
DARK_GRAY  = rgb(0x32, 0x30, 0x30)
AMBER      = rgb(0xD8, 0x83, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers (copied verbatim from create_presentation.py) ─────────────────────

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_bullet_para(tf, text, font_size=16, color=DARK_GRAY,
                    bold=False, level=0, font_name="Segoe UI", space_before=6):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def header_bar(slide, title, subtitle=None, bar_color=BLUE_MID):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=bar_color)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.1), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.8), Inches(12), Inches(0.4),
                    font_size=14, color=rgb(0xDE, 0xEC, 0xF9))


def footer_bar(slide, text="Content Performance Dashboard  |  How It Works"):
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill_color=LIGHT_GRAY)
    add_textbox(slide, text,
                Inches(0.3), Inches(7.12), Inches(12), Inches(0.3),
                font_size=9, color=MID_GRAY)


# ── Data loading ───────────────────────────────────────────────────────────────

ROOT = Path(__file__).parent.parent.resolve()

OLD_COMMIT = "36245ff"   # Method A — HTML-based scores
NEW_COMMIT = "20bcc2f"   # Method B — norm-markdown scores
SCORES_PATH = "data/scores/ai_readiness_scores.json"
C_SCORES_FILE = ROOT / "data" / "scores" / "ai_readiness_scores_raw_md.json"

DIMENSIONS = [
    "heading_hierarchy", "chunk_autonomy", "context_completeness",
    "entity_normalization", "disambiguation",
    "structured_data_utilization",
]
DIM_LABELS = [
    ("Heading Hierarchy",          "Clear H1→H2→H3 for meaningful chunking"),
    ("Chunk Autonomy",             "Each section standalone without context"),
    ("Context Completeness",       "Key terms and prereqs explained inline"),
    ("Entity Normalization",       "Consistent naming for products/features"),
    ("Disambiguation",             "Ambiguous terms and pronouns resolved"),
    ("Structured Data Utilization","Tables and lists over dense prose"),
]


def git_show_json(commit: str, path: str) -> dict:
    result = subprocess.run(
        ["git", "show", f"{commit}:{path}"],
        capture_output=True, encoding="utf-8", check=True,
    )
    return json.loads(result.stdout)


print("Loading Method A scores (HTML) …")
old_data = git_show_json(OLD_COMMIT, SCORES_PATH)

print("Loading Method B scores (norm-md) …")
new_data = git_show_json(NEW_COMMIT, SCORES_PATH)

print("Loading Method C scores (raw-md) …")
if C_SCORES_FILE.exists():
    c_data = json.loads(C_SCORES_FILE.read_text(encoding="utf-8"))
    print(f"  Loaded {len(c_data)} Method C scores")
else:
    c_data = {}
    print("  Not found — comparison slides will show placeholder")

common_ab = set(old_data) & set(new_data)
c_urls = set(c_data) & common_ab
n_ab = len(common_ab)
n_c = len(c_urls)

# Summary stats
avg_a = sum(old_data[u]["total_recommendations"] for u in common_ab) / n_ab
avg_b = sum(new_data[u]["total_recommendations"] for u in common_ab) / n_ab
avg_c = sum(c_data[u]["total_recommendations"] for u in c_urls) / n_c if c_urls else None
avg_b_c = sum(new_data[u]["total_recommendations"] for u in c_urls) / n_c if c_urls else None

bands = ["Low", "Medium", "High"]
migration_ab: dict = defaultdict(lambda: defaultdict(int))
migration_bc: dict = defaultdict(lambda: defaultdict(int))
for u in common_ab:
    migration_ab[old_data[u]["band"]][new_data[u]["band"]] += 1
for u in c_urls:
    migration_bc[new_data[u]["band"]][c_data[u]["band"]] += 1

band_change_ab = sum(1 for u in common_ab if old_data[u]["band"] != new_data[u]["band"])
band_change_bc = sum(1 for u in c_urls if new_data[u]["band"] != c_data[u]["band"]) if c_urls else 0


# ── Presentation setup ─────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
add_rect(s1, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=BLUE_DARK)
add_rect(s1, Inches(0), Inches(0), Inches(0.15), SLIDE_H, fill_color=BLUE_MID)

add_textbox(s1, "Content Performance Dashboard",
            Inches(0.5), Inches(1.5), Inches(12.5), Inches(0.9),
            font_size=40, bold=True, color=WHITE)
add_textbox(s1, "How It Works",
            Inches(0.5), Inches(2.55), Inches(12.5), Inches(0.7),
            font_size=32, bold=False, color=rgb(0xDE, 0xEC, 0xF9))
add_textbox(s1,
            "Pipeline architecture  ·  AI Readiness scoring  ·  Retrievability scoring  ·  Content ingestion comparison",
            Inches(0.5), Inches(3.4), Inches(12), Inches(0.5),
            font_size=16, color=rgb(0xDE, 0xEC, 0xF9))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
add_rect(s2, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s2, "Two gaps in content intelligence",
           "We know how humans engage — we're flying blind on how AI does")
footer_bar(s2)

# Left panel — known
add_rect(s2, Inches(0.3), Inches(1.5), Inches(5.7), Inches(5.3),
         fill_color=rgb(0xE8, 0xF4, 0xE8), line_color=GREEN, line_width=Pt(2))
add_textbox(s2, "✓  Human engagement",
            Inches(0.5), Inches(1.65), Inches(5.3), Inches(0.55),
            font_size=18, bold=True, color=GREEN)
add_textbox(s2, "Power BI tells us:",
            Inches(0.5), Inches(2.3), Inches(5.3), Inches(0.35),
            font_size=13, bold=True, color=DARK_GRAY)
known_box = s2.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(5.3), Inches(3.5))
ktf = known_box.text_frame; ktf.word_wrap = True
for item in ["Page views, bounce rate, CTR",
             "Organic search referrals",
             "Helpfulness ratings",
             "Time on page, scroll depth"]:
    p = ktf.paragraphs[0] if item == "Page views, bounce rate, CTR" else ktf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run(); r.text = f"•  {item}"
    r.font.size = Pt(13); r.font.color.rgb = DARK_GRAY; r.font.name = "Segoe UI"

# Right panel — unknown
add_rect(s2, Inches(7.3), Inches(1.5), Inches(5.7), Inches(5.3),
         fill_color=rgb(0xFD, 0xF0, 0xED), line_color=ORANGE, line_width=Pt(2))
add_textbox(s2, "✗  AI engagement",
            Inches(7.5), Inches(1.65), Inches(5.3), Inches(0.55),
            font_size=18, bold=True, color=ORANGE)
add_textbox(s2, "Power BI does NOT tell us:",
            Inches(7.5), Inches(2.3), Inches(5.3), Inches(0.35),
            font_size=13, bold=True, color=DARK_GRAY)
unknown_box = s2.shapes.add_textbox(Inches(7.5), Inches(2.7), Inches(5.3), Inches(3.5))
utf = unknown_box.text_frame; utf.word_wrap = True
for item in ["Whether Copilot retrieves this article",
             "Whether the article is well-chunked for RAG",
             "Which articles are outcompeted in AI search",
             "How AI-readable the content structure is"]:
    p = utf.paragraphs[0] if item == "Whether Copilot retrieves this article" else utf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run(); r.text = f"•  {item}"
    r.font.size = Pt(13); r.font.color.rgb = DARK_GRAY; r.font.name = "Segoe UI"

# Gap arrow
add_textbox(s2, "?", Inches(6.2), Inches(3.5), Inches(0.9), Inches(0.9),
            font_size=40, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_textbox(s2, "\"When a user asks Copilot a question, does our article show up?  We don't know.\"",
            Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
            font_size=13, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
add_rect(s3, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s3, "End-to-end pipeline", "From engagement export to enriched dashboard")
footer_bar(s3)

arch_boxes = [
    ("Power BI\nCSV",          LIGHT_GRAY, DARK_GRAY, "Engagement metrics\n(URLs, views, etc.)"),
    ("Content\nCache",         BLUE_MID,   WHITE,     "3 ingestion methods\n— see Stage 1"),
    ("AI Readiness\nScorer",   ORANGE,     WHITE,     "10 LLM calls / article\n→ band (High/Med/Low)"),
    ("Retrievability\nScorer", GREEN,      WHITE,     "Question gen → KS query\n→ hit-rate score"),
    ("Enriched\nCSV",          LIGHT_GRAY, DARK_GRAY, "Original cols +\nnew signals"),
    ("Streamlit\nDashboard",   BLUE_DARK,  WHITE,     "Table · Scatter · \nPer-article report"),
]

box_w = Inches(1.85)
box_h = Inches(1.2)
start_x = Inches(0.25)
box_y = Inches(2.1)
gap = Inches(0.2)

for i, (label, fill, text_col, note) in enumerate(arch_boxes):
    bx = start_x + i * (box_w + gap + Inches(0.35))
    add_rect(s3, bx, box_y, box_w, box_h, fill_color=fill)
    add_textbox(s3, label, bx, box_y, box_w, box_h,
                font_size=12, bold=True, color=text_col, align=PP_ALIGN.CENTER)
    add_textbox(s3, note, bx, box_y + Inches(1.3), box_w, Inches(0.8),
                font_size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)
    if i < len(arch_boxes) - 1:
        add_textbox(s3, "→", bx + box_w + Inches(0.05), box_y,
                    Inches(0.3), box_h,
                    font_size=20, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

# Stage labels
stage_info = [
    (Inches(2.25), "Stage 1\nContent\nIngestion"),
    (Inches(4.65), "Stage 2\nAI Readiness\nScoring"),
    (Inches(7.05), "Stage 3\nRetrievability\nScoring"),
]
for sx, slabel in stage_info:
    add_rect(s3, sx, Inches(3.65), Inches(2.0), Inches(0.6),
             fill_color=LIGHT_GRAY, line_color=MID_GRAY, line_width=Pt(0.5))
    add_textbox(s3, slabel, sx, Inches(3.68), Inches(2.0), Inches(0.55),
                font_size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s3, "Input", Inches(0.25), Inches(3.65), Inches(1.85), Inches(0.4),
            font_size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s3, "Output", Inches(9.45), Inches(3.65), Inches(3.85), Inches(0.4),
            font_size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STAGE 1: CONTENT INGESTION
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
add_rect(s4, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s4, "Stage 1: Fetching and caching article content",
           "Three ingestion methods — each gives the LLM a different view of the same article")
footer_bar(s4)

methods = [
    (BLUE_MID, "Method A", "HTML crawl",
     ["Source: live URL (HTTP fetch)",
      "Parser: BeautifulSoup → plain text",
      "Headings: all flattened to same level",
      "Links/images: stripped",
      "Code blocks: stripped",
      "Tables: stripped",
      "Coverage: 273 articles"],
     "What the LLM sees:\n\"Introduction\nCopilot is an AI tool.\nSee also Configuration.\""),
    (GREEN, "Method B", "Local .md → normalized",
     ["Source: local .md file",
      "Parser: normalize_markdown_content()",
      "Headings: flattened to ##",
      "Links/images: stripped",
      "Code blocks: stripped (fence removed)",
      "Tables: separator rows removed",
      "Coverage: 273 articles"],
     "What the LLM sees:\n\"## Introduction\nCopilot is an AI tool.\nSee also Configuration.\""),
    (ORANGE, "Method C", "Local .md → raw",
     ["Source: local .md file",
      "Parser: front matter stripped only",
      "Headings: H1–H6 preserved as-is",
      "Links: [text](url) preserved",
      "Code blocks: ``` preserved",
      "Tables: | pipes | preserved",
      "Coverage: 155 articles (matched .md)"],
     "What the LLM sees:\n\"# Introduction\nCopilot is an [AI tool](url).\n| Feature | Status |\n|---------|--------|\n| Chat    | GA     |\""),
]

col_w = Inches(4.0)
col_top = Inches(1.45)
col_h = Inches(5.4)
for i, (color, label, sublabel, bullets, example) in enumerate(methods):
    cx = Inches(0.2) + i * (col_w + Inches(0.25))
    add_rect(s4, cx, col_top, col_w, col_h,
             fill_color=LIGHT_GRAY, line_color=color, line_width=Pt(2))
    add_rect(s4, cx, col_top, col_w, Inches(0.5), fill_color=color)
    add_textbox(s4, f"{label}  ·  {sublabel}",
                cx + Inches(0.1), col_top + Inches(0.05), col_w - Inches(0.2), Inches(0.4),
                font_size=11, bold=True, color=WHITE)
    bx = s4.shapes.add_textbox(cx + Inches(0.1), col_top + Inches(0.6),
                                col_w - Inches(0.2), Inches(2.6))
    btf = bx.text_frame; btf.word_wrap = True
    for j, b in enumerate(bullets):
        p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = b
        r.font.size = Pt(10); r.font.color.rgb = DARK_GRAY; r.font.name = "Segoe UI"
    add_rect(s4, cx + Inches(0.1), col_top + Inches(3.3),
             col_w - Inches(0.2), Inches(1.8),
             fill_color=WHITE, line_color=color, line_width=Pt(0.75))
    add_textbox(s4, example,
                cx + Inches(0.2), col_top + Inches(3.35),
                col_w - Inches(0.4), Inches(1.7),
                font_size=8, color=DARK_GRAY, font_name="Segoe UI Mono", wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STAGE 2: AI READINESS SCORING
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
add_rect(s5, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s5, "Stage 2: AI Readiness — is this article well-structured for AI?",
           "10 concurrent LLM calls per article, each evaluating a RAG-quality dimension",
           bar_color=ORANGE)
footer_bar(s5)

# Left: 10 dimensions
add_textbox(s5, "10 Dimensions",
            Inches(0.3), Inches(1.45), Inches(6.5), Inches(0.4),
            font_size=13, bold=True, color=ORANGE)
dim_box = s5.shapes.add_textbox(Inches(0.3), Inches(1.9), Inches(6.4), Inches(5.2))
dtf = dim_box.text_frame; dtf.word_wrap = True
for i, (name, desc) in enumerate(DIM_LABELS):
    p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
    p.space_before = Pt(2)
    r = p.add_run(); r.text = f"  {name}"
    r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = ORANGE; r.font.name = "Segoe UI"
    p2 = dtf.add_paragraph()
    p2.space_before = Pt(0)
    r2 = p2.add_run(); r2.text = f"    {desc}"
    r2.font.size = Pt(9); r2.font.color.rgb = MID_GRAY; r2.font.name = "Segoe UI"

# Right: how scoring works
add_textbox(s5, "How It Works",
            Inches(7.1), Inches(1.45), Inches(5.8), Inches(0.4),
            font_size=13, bold=True, color=ORANGE)
steps = [
    (ORANGE, "Per-dimension LLM call",
     "Each of the 10 dimensions runs as a separate, concurrent call to Azure OpenAI. "
     "A shared prompt describes the Knowledge Service: hybrid search, 500-token chunks, "
     "H1+H2 prepended. Only high-confidence issues are flagged — zero recommendations "
     "is valid and expected for a well-structured article."),
    (BLUE_MID, "Recommendation structure",
     "Each recommendation has three fields:\n"
     "  •  Evidence — specific quote showing the issue\n"
     "  •  Action   — imperative fix the editor should make\n"
     "  •  Impact   — the retrieval benefit of fixing it"),
    (GREEN, "Band assignment",
     "Total recommendations across all 10 dimensions → band:\n"
     "  •  High    (0–3 recs)  well-structured for RAG\n"
     "  •  Medium  (4–8 recs)  meaningful gaps; improve\n"
     "  •  Low     (9+ recs)   significant issues; prioritize"),
]
step_top = Inches(1.9)
for color, title, body in steps:
    add_rect(s5, Inches(7.1), step_top, Inches(5.9), Inches(0.35), fill_color=color)
    add_textbox(s5, title, Inches(7.2), step_top + Inches(0.03),
                Inches(5.7), Inches(0.3), font_size=11, bold=True, color=WHITE)
    bb = s5.shapes.add_textbox(Inches(7.1), step_top + Inches(0.38),
                                Inches(5.9), Inches(1.05))
    bbtf = bb.text_frame; bbtf.word_wrap = True
    for j, line in enumerate(body.split("\n")):
        bp = bbtf.paragraphs[0] if j == 0 else bbtf.add_paragraph()
        br = bp.add_run(); br.text = line
        br.font.size = Pt(10); br.font.color.rgb = DARK_GRAY; br.font.name = "Segoe UI"
    step_top += Inches(1.55)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STAGE 3: RETRIEVABILITY SCORING
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
add_rect(s6, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s6, "Stage 3: Retrievability — does AI actually surface this article?",
           "Hit-rate test: does the Knowledge Service return this article for relevant questions?",
           bar_color=GREEN)
footer_bar(s6)

ret_steps = [
    (BLUE_MID, "Step 1\nQuestion Generation",
     ["LLM reads the article and generates 10 questions:",
      "• 5 natural-language questions (vector retrieval)",
      "• 5 BM25 keyword queries (sparse retrieval)",
      "",
      "Grounded in actual article content — not generic."]),
    (rgb(0x10, 0x60, 0x90), "Step 2\nKnowledge Service Query",
     ["Each question sent to Microsoft Learn KS:",
      "• Returns top 5 chunks (hybrid: keyword + vector + reranking)",
      "• Each chunk: up to 500 tokens, H1+H2 prepended",
      "",
      "HIT = article URL appears in any of the 5 chunks."]),
    (GREEN, "Step 3\nScore Calculation",
     ["Score = (hits / 10 questions) × 100",
      "",
      "Example: 7 of 10 → score = 70",
      "",
      "Ranges:",
      "  70–100  Strong retrieval signal",
      "  40–69   Moderate; some questions miss",
      "  0–39    Weak; article is outcompeted"]),
    (ORANGE, "Why This Matters",
     ["High page views + low retrievability = danger zone:",
      "",
      "• Users find the article via direct search",
      "• But asking Copilot the same question returns",
      "  a different (possibly wrong) article",
      "",
      "Invisible without the Retrievability score."]),
]

card_left = Inches(0.2)
for color, title, body_lines in ret_steps:
    add_rect(s6, card_left, Inches(1.45), Inches(3.1), Inches(5.3),
             fill_color=LIGHT_GRAY, line_color=color, line_width=Pt(2))
    add_rect(s6, card_left, Inches(1.45), Inches(3.1), Inches(0.55), fill_color=color)
    add_textbox(s6, title,
                card_left + Inches(0.1), Inches(1.47), Inches(2.9), Inches(0.5),
                font_size=11, bold=True, color=WHITE)
    bb2 = s6.shapes.add_textbox(card_left + Inches(0.1), Inches(2.1),
                                  Inches(2.9), Inches(4.4))
    bb2tf = bb2.text_frame; bb2tf.word_wrap = True
    for j, line in enumerate(body_lines):
        bp2 = bb2tf.paragraphs[0] if j == 0 else bb2tf.add_paragraph()
        bp2.space_before = Pt(2)
        br2 = bp2.add_run(); br2.text = line
        br2.font.size = Pt(10); br2.font.color.rgb = DARK_GRAY; br2.font.name = "Segoe UI"
    card_left += Inches(3.3)

add_rect(s6, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.45),
         fill_color=rgb(0xE8, 0xF4, 0xE8), line_color=GREEN, line_width=Pt(1))
add_textbox(s6, "Score  =  (questions where article appears in top-5 KS chunks)  ÷  10  ×  100",
            Inches(0.5), Inches(6.62), Inches(12.3), Inches(0.38),
            font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STAGE 4: DASHBOARD OUTPUT
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
add_rect(s7, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s7, "Stage 4: What the pipeline produces",
           "Three views in the Streamlit dashboard")
footer_bar(s7)

views = [
    (BLUE_MID, "1  Data Table",
     ["All original Power BI columns intact",
      "New: AI Readiness band (High / Medium / Low)",
      "New: Retrievability score (0–100)",
      "New: Report ↗ link per article",
      "",
      "Filterable, sortable — your Power BI",
      "report with AI performance columns added."]),
    (ORANGE, "2  Portfolio Scatter  ★",
     ["X axis: Retrievability score",
      "Y axis: Page views",
      "Color: AI Readiness band",
      "",
      "Top-left quadrant = high-traffic articles",
      "Copilot fails to surface — your editing backlog.",
      "",
      "Priority table: Low band AND score < 40,",
      "ranked by page views."]),
    (GREEN, "3  Per-Article HTML Report",
     ["Band chip + overall score",
      "Dimension bar chart (recs per dimension)",
      "Recommendation cards:",
      "  •  Evidence — what the issue is",
      "  •  Action   — what to fix",
      "  •  Impact   — retrieval benefit",
      "",
      "Open in VS Code Web — one click to edit."]),
]

for i, (color, title, bullets) in enumerate(views):
    cx = Inches(0.3) + i * Inches(4.35)
    add_rect(s7, cx, Inches(1.5), Inches(4.1), Inches(5.3),
             fill_color=LIGHT_GRAY, line_color=color, line_width=Pt(2))
    add_rect(s7, cx, Inches(1.5), Inches(4.1), Inches(0.5), fill_color=color)
    add_textbox(s7, title, cx + Inches(0.1), Inches(1.52),
                Inches(3.9), Inches(0.45), font_size=13, bold=True, color=WHITE)
    vb = s7.shapes.add_textbox(cx + Inches(0.15), Inches(2.1),
                                Inches(3.8), Inches(4.4))
    vtf = vb.text_frame; vtf.word_wrap = True
    for j, b in enumerate(bullets):
        p = vtf.paragraphs[0] if j == 0 else vtf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run(); r.text = b
        r.font.size = Pt(11); r.font.color.rgb = DARK_GRAY; r.font.name = "Segoe UI"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY THREE INGESTION METHODS (motivation)
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
add_rect(s8, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s8, "Hypothesis: does richer markdown structure improve LLM assessment?",
           "The same scoring pipeline, three different content inputs")
footer_bar(s8)

# Three-panel: what → hypothesis → how we tested
panels = [
    (BLUE_MID, "What the LLM uses",
     "The scoring prompt feeds the article content directly to the LLM. "
     "The content is the only variable — same 10 dimensions, same prompts, "
     "same Azure OpenAI deployment across all three methods."),
    (ORANGE, "The hypothesis",
     "Preserving heading hierarchy, live links, and table structure gives the LLM "
     "richer structural signals. Richer signals → more accurate issue detection → "
     "fewer spurious recommendations on well-structured content."),
    (GREEN, "How we tested",
     "273 articles scored with Method A (HTML crawl) and Method B (normalized markdown). "
     "155 articles also scored with Method C (raw markdown). "
     "compare_scores.py computes per-article and per-dimension deltas across all three runs."),
]
for i, (color, title, body) in enumerate(panels):
    px = Inches(0.3) + i * Inches(4.35)
    add_rect(s8, px, Inches(1.5), Inches(4.1), Inches(4.0),
             fill_color=LIGHT_GRAY, line_color=color, line_width=Pt(2))
    add_rect(s8, px, Inches(1.5), Inches(4.1), Inches(0.45), fill_color=color)
    add_textbox(s8, title, px + Inches(0.1), Inches(1.52),
                Inches(3.9), Inches(0.4), font_size=13, bold=True, color=WHITE)
    add_textbox(s8, body, px + Inches(0.15), Inches(2.05),
                Inches(3.8), Inches(3.2),
                font_size=12, color=DARK_GRAY, wrap=True)

# Comparison table
add_textbox(s8, "Coverage summary",
            Inches(0.3), Inches(5.7), Inches(12.5), Inches(0.35),
            font_size=12, bold=True, color=DARK_GRAY)
tbl_data = [
    ("Method", "Source", "Heading levels", "Links / tables", "Articles"),
    ("A — HTML",    "HTTP crawl",   "flattened",       "stripped",   "273"),
    ("B — Norm-md", "Local .md",    "flattened to ##", "stripped",   "273"),
    ("C — Raw-md",  "Local .md",    "H1–H6 preserved", "preserved",  "155"),
]
col_widths = [Inches(1.8), Inches(1.8), Inches(2.5), Inches(2.5), Inches(1.5)]
col_starts = [Inches(0.3)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + Inches(0.05))
for row_i, row in enumerate(tbl_data):
    bg = rgb(0xE0, 0xEC, 0xF8) if row_i == 0 else (LIGHT_GRAY if row_i % 2 == 1 else WHITE)
    row_top = Inches(6.1) + row_i * Inches(0.3)
    for ci, (cell, cw, cx_) in enumerate(zip(row, col_widths, col_starts)):
        add_rect(s8, cx_, row_top, cw, Inches(0.28), fill_color=bg)
        add_textbox(s8, cell, cx_ + Inches(0.05), row_top + Inches(0.02),
                    cw - Inches(0.1), Inches(0.25),
                    font_size=9, bold=(row_i == 0), color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — METHOD A vs B vs C: WHAT THE LLM SEES
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
add_rect(s9, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s9, "What each method feeds to the LLM",
           "Same article — three different representations")
footer_bar(s9)

method_cards = [
    (BLUE_MID, "Method A — HTML crawl",
     "Introduction\nCopilot is an AI tool that helps\nyou work more efficiently.\nSee also Configuration.\n\n"
     "Features\nCopilot supports Word, Excel,\nPowerPoint, Teams, and Outlook.",
     ["Structure: none", "Headings: lost", "Links: stripped", "Tables: stripped", "Code: stripped"]),
    (GREEN, "Method B — Normalized markdown",
     "## Introduction\nCopilot is an AI tool that helps\nyou work more efficiently.\nSee also Configuration.\n\n"
     "## Features\nCopilot supports Word, Excel,\nPowerPoint, Teams, and Outlook.",
     ["Structure: ## only", "Headings: flattened", "Links: stripped", "Tables: partial", "Code: text only"]),
    (ORANGE, "Method C — Raw markdown",
     "# Introduction\nCopilot is an [AI tool](url) that\nhelps you work more efficiently.\nSee also [Configuration](url).\n\n"
     "## Features\n| App | Status |\n|-----|--------|\n| Word | GA |\n| Teams | GA |",
     ["Structure: full H1–H6", "Headings: preserved", "Links: [text](url)", "Tables: | pipes |", "Code: ``` fences"]),
]
for i, (color, title, example, props) in enumerate(method_cards):
    cx = Inches(0.2) + i * Inches(4.35)
    add_rect(s9, cx, Inches(1.5), Inches(4.1), Inches(5.1),
             fill_color=LIGHT_GRAY, line_color=color, line_width=Pt(2))
    add_rect(s9, cx, Inches(1.5), Inches(4.1), Inches(0.45), fill_color=color)
    add_textbox(s9, title, cx + Inches(0.1), Inches(1.52),
                Inches(3.9), Inches(0.4), font_size=11, bold=True, color=WHITE)
    # Code example box
    add_rect(s9, cx + Inches(0.1), Inches(2.05), Inches(3.9), Inches(2.4),
             fill_color=WHITE, line_color=color, line_width=Pt(0.75))
    add_textbox(s9, example, cx + Inches(0.2), Inches(2.1),
                Inches(3.7), Inches(2.3), font_size=8, color=DARK_GRAY,
                font_name="Segoe UI Mono", wrap=True)
    # Properties
    pb = s9.shapes.add_textbox(cx + Inches(0.15), Inches(4.55), Inches(3.8), Inches(1.8))
    ptf = pb.text_frame; ptf.word_wrap = True
    for j, prop in enumerate(props):
        p = ptf.paragraphs[0] if j == 0 else ptf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = f"•  {prop}"
        r.font.size = Pt(10); r.font.color.rgb = DARK_GRAY; r.font.name = "Segoe UI"

# Key insight
add_rect(s9, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.5),
         fill_color=rgb(0xFF, 0xF4, 0xE0), line_color=AMBER, line_width=Pt(1.5))
add_textbox(s9,
            "B and C use the same source files — score differences are purely structural. "
            "A uses HTTP-crawled HTML which may differ slightly from the .md source.",
            Inches(0.5), Inches(6.58), Inches(12.3), Inches(0.4),
            font_size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SCORE COMPARISON RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
add_rect(s10, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=WHITE)
header_bar(s10, "Score comparison: does structure preservation change results?",
           f"Method A vs B: {n_ab} articles  ·  Method C: {n_c} articles matched")
footer_bar(s10)

if c_data:
    # ── Summary stats table ────────────────────────────────────────────────────
    add_textbox(s10, "Average total recommendations per method",
                Inches(0.3), Inches(1.5), Inches(6.0), Inches(0.35),
                font_size=12, bold=True, color=DARK_GRAY)

    tbl_rows = [
        ("Method", "N", "Avg recs", "vs A", "vs B"),
        (f"A — HTML ({OLD_COMMIT})", str(n_ab), f"{avg_a:.1f}", "—", "—"),
        (f"B — Norm-md ({NEW_COMMIT})", str(n_ab), f"{avg_b:.1f}", f"{avg_b - avg_a:+.1f}", "—"),
        (f"C — Raw-md", str(n_c), f"{avg_c:.1f}",
         f"{avg_c - avg_a:+.1f}", f"{avg_c - avg_b_c:+.1f}"),
    ]
    cw2 = [Inches(2.8), Inches(0.7), Inches(1.0), Inches(0.8), Inches(0.8)]
    cs2 = [Inches(0.3)]
    for w in cw2[:-1]:
        cs2.append(cs2[-1] + w + Inches(0.05))
    for ri, row in enumerate(tbl_rows):
        bg2 = rgb(0xE0, 0xEC, 0xF8) if ri == 0 else (LIGHT_GRAY if ri % 2 == 1 else WHITE)
        rt = Inches(1.95) + ri * Inches(0.38)
        for cell, cw_, cx_ in zip(row, cw2, cs2):
            add_rect(s10, cx_, rt, cw_, Inches(0.35), fill_color=bg2)
            is_delta = cell.startswith("+") or (cell.startswith("-") and cell != "—")
            cell_color = GREEN if (cell.startswith("-") and cell != "—") else \
                         ORANGE if cell.startswith("+") else DARK_GRAY
            add_textbox(s10, cell, cx_ + Inches(0.05), rt + Inches(0.04),
                        cw_ - Inches(0.1), Inches(0.28),
                        font_size=9, bold=(ri == 0), color=cell_color)

    # ── Band change counts ─────────────────────────────────────────────────────
    add_textbox(s10, "Band changes",
                Inches(0.3), Inches(3.65), Inches(5.5), Inches(0.35),
                font_size=12, bold=True, color=DARK_GRAY)
    bc_rows = [
        ("Comparison", "Band changes", "% of N"),
        (f"A → B  (N={n_ab})",
         str(band_change_ab), f"{100*band_change_ab/n_ab:.0f}%"),
        (f"B → C  (N={n_c})",
         str(band_change_bc), f"{100*band_change_bc/n_c:.0f}%" if n_c else "—"),
    ]
    bcw = [Inches(2.5), Inches(1.5), Inches(1.2)]
    bcs = [Inches(0.3)]
    for w in bcw[:-1]:
        bcs.append(bcs[-1] + w + Inches(0.05))
    for ri, row in enumerate(bc_rows):
        bg3 = rgb(0xE0, 0xEC, 0xF8) if ri == 0 else (LIGHT_GRAY if ri % 2 == 1 else WHITE)
        rt2 = Inches(4.05) + ri * Inches(0.35)
        for cell, cw_, cx_ in zip(row, bcw, bcs):
            add_rect(s10, cx_, rt2, cw_, Inches(0.32), fill_color=bg3)
            add_textbox(s10, cell, cx_ + Inches(0.05), rt2 + Inches(0.03),
                        cw_ - Inches(0.1), Inches(0.26),
                        font_size=9, bold=(ri == 0), color=DARK_GRAY)

    # ── B→C migration matrix ───────────────────────────────────────────────────
    add_textbox(s10, "B → C band migration matrix",
                Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.35),
                font_size=12, bold=True, color=DARK_GRAY)
    m_top = Inches(1.95)
    m_bands = ["Low", "Medium", "High"]
    m_headers = ["B \\ C"] + m_bands + ["Row Total"]
    m_cw = [Inches(1.3), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)]
    m_cs = [Inches(7.0)]
    for w in m_cw[:-1]:
        m_cs.append(m_cs[-1] + w + Inches(0.05))
    # Header row
    for cell, cw_, cx_ in zip(m_headers, m_cw, m_cs):
        add_rect(s10, cx_, m_top, cw_, Inches(0.32), fill_color=rgb(0xE0, 0xEC, 0xF8))
        add_textbox(s10, cell, cx_ + Inches(0.05), m_top + Inches(0.03),
                    cw_ - Inches(0.1), Inches(0.26), font_size=9, bold=True, color=DARK_GRAY)
    for ri, b_band in enumerate(m_bands):
        rt3 = m_top + (ri + 1) * Inches(0.37)
        row_vals = [str(migration_bc[b_band][c_band]) for c_band in m_bands]
        row_total = sum(migration_bc[b_band][c_band] for c_band in m_bands)
        row_data = [f"  {b_band}"] + row_vals + [str(row_total)]
        bg4 = LIGHT_GRAY if ri % 2 == 0 else WHITE
        for cell, cw_, cx_ in zip(row_data, m_cw, m_cs):
            add_rect(s10, cx_, rt3, cw_, Inches(0.34), fill_color=bg4)
            add_textbox(s10, cell, cx_ + Inches(0.05), rt3 + Inches(0.04),
                        cw_ - Inches(0.1), Inches(0.28), font_size=9, color=DARK_GRAY)

    # Insight line
    avg_delta_bc = avg_c - avg_b_c
    direction = "fewer" if avg_delta_bc < 0 else "more"
    insight = (
        f"Method C produces on average {abs(avg_delta_bc):.1f} {direction} recommendations than B on matched articles. "
        f"B→C band changes: {band_change_bc} / {n_c} articles ({100*band_change_bc/n_c:.0f}%)."
        if n_c else ""
    )
    add_rect(s10, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.55),
             fill_color=rgb(0xE8, 0xF4, 0xE8), line_color=GREEN, line_width=Pt(1))
    add_textbox(s10, insight,
                Inches(0.5), Inches(6.53), Inches(12.3), Inches(0.45),
                font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

else:
    # Placeholder
    add_textbox(s10, "Method C scores not yet available",
                Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
                font_size=24, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_textbox(s10,
                "Run the scoring pipeline to populate this slide:\n\n"
                "python pipeline/run_ai_readiness.py \\\n"
                "    --input data/LMC_Monthly_Engagement_Metrics.csv \\\n"
                "    --cache-dir data/cache_raw_md \\\n"
                "    --output data/scores/ai_readiness_scores_raw_md.json\n\n"
                f"155 cache files are ready in data/cache_raw_md/\n"
                "Requires: AZURE_OPENAI_ENDPOINT + AZURE_OPENAI_KEY in .env",
                Inches(2.5), Inches(2.8), Inches(8.3), Inches(3.5),
                font_size=13, color=DARK_GRAY, font_name="Segoe UI Mono", wrap=True)
    # A vs B stats still available
    add_textbox(s10, f"Method A vs B preview (N={n_ab} articles):",
                Inches(0.3), Inches(5.7), Inches(12.5), Inches(0.35),
                font_size=12, bold=True, color=DARK_GRAY)
    ab_preview = (
        f"Avg recs — A: {avg_a:.1f}  B: {avg_b:.1f}  Δ(B−A): {avg_b - avg_a:+.1f}   |   "
        f"Band changes A→B: {band_change_ab} / {n_ab} ({100*band_change_ab/n_ab:.0f}%)"
    )
    add_textbox(s10, ab_preview,
                Inches(0.3), Inches(6.1), Inches(12.5), Inches(0.4),
                font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────

OUT = ROOT / "docs" / "content-dashboard-how-it-works.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
